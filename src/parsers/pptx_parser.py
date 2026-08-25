"""PPTX 파서: 슬라이드 단위로 텍스트를 추출한다.

회의록이 전부 이 형식이다. 슬라이드 번호가 출처 위치가 되며,
발표자 노트는 별도 블록으로 남긴다(본문과 성격이 다르다).
"""

from pathlib import Path

from pptx import Presentation

try:
    from .base import Block
except ImportError:
    from base import Block


def _table_text(shape) -> str:
    rows = []
    for row in shape.table.rows:
        cells = [c.text.strip().replace("\n", " ") for c in row.cells]
        if any(cells):
            rows.append(" | ".join(cells))
    return "\n".join(rows)


def _shape_text(shape) -> str:
    if shape.has_text_frame:
        return shape.text_frame.text.strip()
    return ""


def parse_pptx(path: str | Path) -> list[Block]:
    prs = Presentation(path)
    blocks: list[Block] = []

    for slide_no, slide in enumerate(prs.slides, start=1):
        loc = f"슬라이드 {slide_no}"

        # 표는 별도 블록으로 둔다. 성과 표를 정형 데이터로 뽑아낼 수 있어야 한다.
        for shape in slide.shapes:
            if shape.has_table and (t := _table_text(shape)):
                blocks.append(Block(len(blocks) + 1, "table", t, locator=loc))

        parts = [t for shape in slide.shapes if not shape.has_table and (t := _shape_text(shape))]
        if parts:
            blocks.append(Block(len(blocks) + 1, "body", "\n".join(parts), locator=loc))

        if slide.has_notes_slide:
            # 노트 슬라이드가 있어도 텍스트 프레임이 비어 있을 수 있다
            frame = slide.notes_slide.notes_text_frame
            notes = frame.text.strip() if frame is not None and frame.text else ""
            if notes:
                blocks.append(
                    Block(len(blocks) + 1, "body", notes, locator=f"{loc} (발표자 노트)")
                )

    return blocks


if __name__ == "__main__":
    import sys

    target = Path(sys.argv[1])
    blocks = parse_pptx(target)

    print(f"파일        : {target.name}")
    print(f"텍스트 추출 : {len(blocks)} 블록, 총 {sum(len(b.text) for b in blocks):,}자")
    print("-" * 60)
    for b in blocks[:3]:
        print(f"[{b.locator}] {b.text[:400]}")
        print("-" * 60)
